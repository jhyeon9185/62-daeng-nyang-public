
import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- Constants & Configuration ---
MD_FILE_PATH = r'c:\workspace\62dn\.notes\lecture_ppt_and_script.md'
OUTPUT_PPTX_PATH = r'c:\workspace\62dn\.notes\62dn_lecture_modern.pptx'

# Colors (Apple Style V5)
BG_COLOR = RGBColor(0, 0, 0)         # Pull Black
TEXT_TITLE = RGBColor(255, 255, 255) # Pure White
TEXT_BODY = RGBColor(161, 161, 166)  # Apple Gray (#A1A1A6)
ACCENT_COLOR = RGBColor(10, 132, 255) # System Blue (iOS) - for subtle accents

# Fonts
FONT_NAME_BOLD = 'Malgun Gothic' 
FONT_NAME_REGULAR = 'Malgun Gothic'

def apply_background(slide):
    """Applies a solid dark background to the slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def get_optimized_font_size(text, max_char_count=20, default_size=60, min_size=40):
    """Calculates font size based on text length."""
    length = len(text)
    if length <= max_char_count:
        return Pt(default_size)
    
    # Simple linear reduction
    excess = length - max_char_count
    reduction = excess * 1.5 
    new_size = default_size - reduction
    return Pt(max(new_size, min_size))

def add_title_slide(prs, title_text, subtitle_text):
    """Creates a minimalist Apple-style title slide."""
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide)

    # 1. Title (Huge, White, Centered)
    # Start slightly higher to allow for big impact
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(3))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    
    # Dynamic Font Sizing (Target 60pt+)
    p.font.size = get_optimized_font_size(title_text, max_char_count=12, default_size=64, min_size=44)
    p.font.name = FONT_NAME_BOLD
    p.font.bold = True
    p.font.color.rgb = TEXT_TITLE
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.1

    # 2. Subtitle (Gray, Elegant)
    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(7), Inches(1.2))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle_text
        p_sub.font.size = Pt(28)
        p_sub.font.name = FONT_NAME_REGULAR
        p_sub.font.color.rgb = TEXT_BODY
        p_sub.alignment = PP_ALIGN.CENTER
        p_sub.line_spacing = 1.3

    return slide


def add_content_slide(prs, title_text, content_items, visual_suggestion=""):
    """Creates a content slide. Split layout if visual exists, otherwise centered list."""
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide)

    # Layout Strategy:
    # If Visual Suggestion exists -> Split Screen (Text Left 50%, Visual Right 50%)
    # If No Visual -> Text Centered or Wide Left with breathing room

    has_visual = bool(visual_suggestion)
    
    # --- Title ---
    # Position: Top Left
    title_font_size = get_optimized_font_size(title_text, max_char_count=20, default_size=48, min_size=36)
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(9), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = title_font_size
    p.font.name = FONT_NAME_BOLD
    p.font.bold = True
    p.font.color.rgb = TEXT_TITLE
    p.alignment = PP_ALIGN.LEFT

    # --- Content ---
    if has_visual:
        # Split Layout: Text on Left
        content_left = Inches(0.8)
        content_top = Inches(1.8)
        content_width = Inches(4.5)
        content_height = Inches(3.5)
        
        # Visual Placeholder on Right
        vis_left = Inches(5.5)
        vis_top = Inches(1.8)
        vis_width = Inches(4.0)
        vis_height = Inches(3.2)
        
        # Draw Placeholder for Visual
        vis_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, vis_left, vis_top, vis_width, vis_height
        )
        vis_shape.fill.solid()
        vis_shape.fill.fore_color.rgb = RGBColor(28, 28, 30) # Dark Gray Placeholder
        vis_shape.line.fill.background()
        
        # Add "Image Placeholder" text
        tf_vis = vis_shape.text_frame
        p_vis = tf_vis.paragraphs[0]
        p_vis.text = "Visual Asset"
        p_vis.font.color.rgb = TEXT_BODY
        p_vis.alignment = PP_ALIGN.CENTER

    else:
        # Standard Layout: Wide Text, slightly indented
        content_left = Inches(0.8)
        content_top = Inches(1.8)
        content_width = Inches(8.4)
        content_height = Inches(3.5)

    # Draw Content Text
    text_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    tf_content = text_box.text_frame
    tf_content.word_wrap = True
    
    # Dynamic content font size
    base_font_size = 32 # Large, readable
    if len(content_items) > 5:
        base_font_size = 28
    if len(content_items) > 8:
        base_font_size = 24
    if has_visual and len(content_items) > 5:
        base_font_size = 24 # Reduce more for split screen

    for item in content_items:
        p = tf_content.add_paragraph()
        
        clean_text = item.strip()
        level = 0
        if item.startswith("    "):
            level = 1
            clean_text = item.strip()
        
        p.text = clean_text # No bullets, just clean text? Or maybe simple bullet
        
        # Add a custom bullet manually if it's a list item
        # Apple style often uses no bullets for big statements, but for lists we need them.
        # Let's add a simple "• " prefix if it's not already there
        if not p.text.startswith("•") and not p.text[0].isdigit():
             p.text = "• " + p.text

        p.level = level
        p.font.size = Pt(base_font_size)
        p.font.name = FONT_NAME_REGULAR
        p.font.color.rgb = TEXT_BODY
        p.space_after = Pt(20) # Energetic spacing
        p.line_spacing = 1.3 

    return slide

def parse_markdown(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # Normalize newlines
    content = content.replace('\r\n', '\n')
    
    # Split by slide separator
    raw_slides = content.split('\n---\n')
    
    parsed_slides = []
    
    for i, raw_text in enumerate(raw_slides):
        raw_text = raw_text.strip()
        if not raw_text:
            continue
            
        slide_data = {
            'type': 'content', # default
            'title': '',
            'subtitle': '',
            'content': [],
            'visual': '',
            'script': ''
        }
        
        lines = raw_text.split('\n')
        current_section = None
        
        # Heuristic for Title Slide
        if 'Slide 1]' in raw_text or (i == 0 and 'Slide' not in raw_text):
            slide_data['type'] = 'title'

        for line in lines:
            line = line.strip()
            if not line: continue
            
            # Skip slide markers
            if line.startswith('## [Slide'):
                continue
            
            # Section Headers detection
            if line.startswith('### 제목'):
                current_section = 'title'
                continue
            elif line.startswith('### 핵심 내용') or line.startswith('### 내용'):
                current_section = 'content'
                continue
            elif line.startswith('### 시각 자료 제안'):
                current_section = 'visual'
                continue
            elif line.startswith('### 대본') or line.startswith('### 스크립트'):
                current_section = 'script'
                continue
            
            # Parse Content based on section
            if current_section == 'title':
                if line.startswith('**'):
                    slide_data['title'] = line.replace('**', '')
                elif line.startswith('부제:'):
                    slide_data['subtitle'] = line.replace('부제:', '').strip()
                elif not slide_data['title']: # Fallback if no **
                    slide_data['title'] = line
                    
            elif current_section == 'content':
                # Remove Markdown list markers
                if line.startswith('- ') or line.startswith('* '):
                    slide_data['content'].append(line[2:])
                elif line.startswith('    - '):
                    slide_data['content'].append("    " + line.strip('- ').strip())
                elif line[0].isdigit() and '. ' in line[:4]: # "1. ", "2. " etc
                    slide_data['content'].append(line)
                elif line.startswith('>'): # Blockquotes
                    slide_data['content'].append(line.replace('>', '').strip())
                else:
                    # just append text if it looks like content
                    slide_data['content'].append(line)

            elif current_section == 'visual':
                slide_data['visual'] += line + "\n"

            elif current_section == 'script':
                slide_data['script'] += line + " "

        parsed_slides.append(slide_data)
        
    return parsed_slides

def main():
    if not os.path.exists(MD_FILE_PATH):
        print(f"Error: File not found at {MD_FILE_PATH}")
        return

    print("Parsing Markdown...")
    slides_data = parse_markdown(MD_FILE_PATH)
    print(f"Found {len(slides_data)} slides.")

    prs = Presentation()
    
    # Set slide size to 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    print("Generating PPTX...")
    for idx, slide_info in enumerate(slides_data):
        slide = None
        
        if slide_info['type'] == 'title':
            slide = add_title_slide(prs, slide_info['title'], slide_info['subtitle'])
        else:
            slide = add_content_slide(
                prs, 
                slide_info['title'], 
                slide_info['content'], 
                slide_info['visual'].strip()
            )
        
        # Add Notes (Script + Visual Suggestion)
        notes_text = ""
        if slide_info['visual']:
            notes_text += f"[Visual Suggestion]\n{slide_info['visual'].strip()}\n\n"
        
        if slide_info['script']:
            notes_text += f"[Script]\n{slide_info['script'].strip()}"

        if notes_text:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.text = notes_text

    output_v5 = OUTPUT_PPTX_PATH.replace('.pptx', '_v5.pptx')
    prs.save(output_v5)
    print(f"Successfully generated: {output_v5}")

if __name__ == '__main__':
    main()

