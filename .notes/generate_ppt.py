import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation(md_file_path, output_pptx_path):
    prs = Presentation()

    with open(md_file_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # Split by '---' to separate slides
    # The first chunk might be the overall title/intro if it doesn't start with '---'
    # But based on the file format, the first slide starts after the first '---' or at the beginning.
    # Let's split and iterate.
    
    # Normalize line endings
    content = content.replace('\r\n', '\n')
    
    slides_content = content.split('\n---\n')

    for i, slide_text in enumerate(slides_content):
        slide_text = slide_text.strip()
        if not slide_text:
            continue
            
        # Parse slide content
        title = ""
        subtitle = ""
        bullets = []
        script = ""
        
        lines = slide_text.split('\n')
        
        current_section = None
        
        # Check if it's the title slide (Slide 1 usually)
        is_title_slide = 'Slide 1]' in slide_text or (i == 0 and '# ' in slide_text and 'Slide' not in slide_text)

        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            if line.startswith('## [Slide'):
                continue # Skip the slide identifier line
            
            if line.startswith('### 제목'):
                current_section = 'title'
                continue
            elif line.startswith('### 핵심 내용') or line.startswith('### 내용'):
                current_section = 'content'
                continue
            elif line.startswith('### 시각 자료 제안'):
                current_section = 'visual' # We might skip this for the PPT text, or add as note
                continue
            elif line.startswith('### 대본') or line.startswith('### 스크립트'):
                current_section = 'script'
                continue
                
            # content processing based on section
            if current_section == 'title':
                if line.startswith('**'):
                    title = line.replace('**', '')
                elif line.startswith('부제:'):
                    subtitle = line.replace('부제:', '').strip()
                elif not title:
                    title = line
                    
            elif current_section == 'content':
                if line.startswith('- ') or line.startswith('* '):
                    bullets.append(line[2:])
                elif line.startswith('1. ') or line.startswith('2. '): # numbered list
                     bullets.append(line)
                elif line.startswith('    - '): # indented bullet
                     bullets.append("    " + line.strip('- ').strip())

            elif current_section == 'script':
                script += line + " "

        # Create Slide
        if is_title_slide and i < 2: # heuristic for title slide
            slide_layout = prs.slide_layouts[0] # Title Slide
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            if hasattr(slide.placeholders[1], 'text'):
                slide.placeholders[1].text = subtitle
        else:
            slide_layout = prs.slide_layouts[1] # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            
            # Set Title
            if slide.shapes.title:
                slide.shapes.title.text = title
            
            # Set Content (Bullets)
            if slide.placeholders[1]:
                tf = slide.placeholders[1].text_frame
                tf.text = "" # clear default
                
                for bullet in bullets:
                    p = tf.add_paragraph()
                    
                    # Handle indentation (simple generic handling)
                    clean_bullet = bullet
                    level = 0
                    if bullet.startswith("    "):
                        level = 1
                        clean_bullet = bullet.strip()
                    
                    p.text = clean_bullet
                    p.level = level

        # Add Script to Notes
        if script:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = script.strip()

    prs.save(output_pptx_path)
    print(f"PPTX generated at: {output_pptx_path}")

if __name__ == "__main__":
    try:
        import pptx
    except ImportError:
        import os
        print("Installing python-pptx...")
        os.system("pip install python-pptx")
        
    create_presentation(r'c:\workspace\62dn\.notes\lecture_ppt_and_script.md', r'c:\workspace\62dn\.notes\62dn_lecture.pptx')
